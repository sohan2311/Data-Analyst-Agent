# Universal Data Analyst Agent
# pip install pandas numpy requests matplotlib seaborn openpyxl python-docx PyPDF2 pillow pdfplumber plotly wordcloud scikit-learn nltk python-pptx pytesseract
# pip install plotly wordcloud scikit-learn nltk python-pptx


import os
import json
import pandas as pd
import numpy as np
from typing import Optional, Dict, Any, List
import requests
from datetime import datetime
import matplotlib.pyplot as plt
import seaborn as sns
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from io import StringIO, BytesIO
import warnings
warnings.filterwarnings('ignore')

# Document processing imports
try:
    from docx import Document
    import PyPDF2
    import pdfplumber
    from PIL import Image
    import pytesseract
    from wordcloud import WordCloud
    import nltk
    from nltk.corpus import stopwords
    from nltk.tokenize import word_tokenize
    from nltk.sentiment import SentimentIntensityAnalyzer
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.cluster import KMeans
    import pptx
    from pptx import Presentation
except ImportError as e:
    print(f"Some optional dependencies are missing: {e}")
    print("Install with: pip install python-docx PyPDF2 pillow pytesseract wordcloud scikit-learn nltk python-pptx pdfplumber")

class UniversalDataAnalystAgent:
    """AI-powered Universal Data Analyst Agent supporting all file types"""
    
    def __init__(self, api_key: str):
        """Initialize the Universal Data Analyst Agent"""
        self.api_key = api_key
        self.base_url = "https://api.together.xyz/v1/chat/completions"
        self.model = "meta-llama/Llama-4-Maverick-17B-128E-Instruct-FP8"
        self.dataset = None
        self.text_content = ""
        self.file_type = ""
        self.file_info = {}
        self.analysis_history = []
        self.visualizations = []
    
        self.headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    
    # Initializing NLTK data with SSL
        self._initialize_nltk()
    
    
    def _initialize_nltk(self):
        """Initialize NLTK data with SSL certificate handling"""
    try:  # This line needs proper indentation
        import ssl
        

        try:
            _create_unverified_https_context = ssl._create_unverified_context
        except AttributeError:
            pass
        else:
            ssl._create_default_https_context = _create_unverified_https_context
        

        nltk.download('punkt', quiet=True)
        nltk.download('stopwords', quiet=True)
        nltk.download('vader_lexicon', quiet=True)
        print(" NLTK data initialized successfully")
        
    except Exception as e:
        print(f" NLTK initialization warning: {str(e)}")
        print("Text analysis features may be limited, but basic functionality will work.")



    
    def query_llm(self, prompt: str, context: str = "") -> str:
        """Query the Together.ai LLM with given prompt and context"""
        try:
            full_prompt = f"Context: {context}\n\nQuery: {prompt}" if context else prompt
            
            payload = {
                "model": self.model,
                "messages": [
                    {
                        "role": "system",
                        "content": "You are an expert data analyst capable of analyzing any type of data - structured data (CSV, Excel), documents (PDF, Word, PowerPoint), text files, and images. Provide clear, accurate, and actionable insights. Format responses with sections and bullet points for readability."
                    },
                    {
                        "role": "user",
                        "content": full_prompt
                    }
                ],
                "max_tokens": 3000,
                "temperature": 0.7,
                "top_p": 0.9
            }
            
            response = requests.post(self.base_url, headers=self.headers, json=payload)
            
            if response.status_code == 200:
                result = response.json()
                return result['choices'][0]['message']['content']
            else:
                return f"Error: API request failed with status {response.status_code}: {response.text}"
                
        except Exception as e:
            return f"Error: {str(e)}"
    
    def extract_text_from_docx(self, filename: str) -> str:
        """Extract text from Word documents"""
        try:
            doc = Document(filename)
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            return '\n'.join(text)
        except Exception as e:
            return f"Error reading DOCX: {str(e)}"
    
    def extract_text_from_pdf(self, filename: str) -> str:
        """Extract text from PDF files"""
        try:
            text = ""

            try:
                with pdfplumber.open(filename) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                if text.strip():
                    return text
            except:
                pass
            

            with open(filename, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            return text
        except Exception as e:
            return f"Error reading PDF: {str(e)}"
    
    def extract_text_from_pptx(self, filename: str) -> str:
        """Extract text from PowerPoint presentations"""
        try:
            prs = Presentation(filename)
            text = []
            slide_num = 1
            
            for slide in prs.slides:
                slide_text = [f"\n--- Slide {slide_num} ---"]
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                text.extend(slide_text)
                slide_num += 1
            
            return '\n'.join(text)
        except Exception as e:
            return f"Error reading PPTX: {str(e)}"
    
    def extract_text_from_image(self, filename: str) -> str:
        """Extract text from images using OCR"""
        try:
            image = Image.open(filename)
            text = pytesseract.image_to_string(image)
            return text
        except Exception as e:
            return f"Error extracting text from image: {str(e)}"
    
    def analyze_image_content(self, filename: str) -> Dict[str, Any]:
        """Analyze image properties and content"""
        try:
            image = Image.open(filename)
            
            analysis = {
                'format': image.format,
                'mode': image.mode,
                'size': image.size,
                'width': image.width,
                'height': image.height,
                'has_transparency': image.mode in ('RGBA', 'LA'),
                'file_size_mb': os.path.getsize(filename) / (1024 * 1024)
            }
            

            try:
                text = pytesseract.image_to_string(image)
                analysis['extracted_text'] = text
                analysis['has_text'] = len(text.strip()) > 10
            except:
                analysis['extracted_text'] = ""
                analysis['has_text'] = False
            
            return analysis
        except Exception as e:
            return {'error': str(e)}
    
    def load_file(self, filename: str) -> str:
        """Universal file loader supporting all file types"""
        try:
            if not os.path.exists(filename):
                return f" File '{filename}' not found."
            
            file_extension = filename.lower().split('.')[-1]
            file_size_mb = os.path.getsize(filename) / (1024 * 1024)
            

            self.dataset = None
            self.text_content = ""
            self.file_type = file_extension
            

            if file_extension == 'csv':
                self.dataset = pd.read_csv(filename)
                self._set_structured_file_info(filename, file_size_mb)
                return self._get_structured_load_message(filename)
                
            elif file_extension in ['xlsx', 'xls']:
                self.dataset = pd.read_excel(filename)
                self._set_structured_file_info(filename, file_size_mb)
                return self._get_structured_load_message(filename)
                
            elif file_extension == 'json':
                with open(filename, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                if isinstance(data, list):
                    self.dataset = pd.DataFrame(data)
                    self._set_structured_file_info(filename, file_size_mb)
                    return self._get_structured_load_message(filename)
                else:
                    self.text_content = json.dumps(data, indent=2)
                    self._set_text_file_info(filename, file_size_mb)
                    return self._get_text_load_message(filename)
            

            elif file_extension == 'txt':
                with open(filename, 'r', encoding='utf-8') as f:
                    self.text_content = f.read()
                self._set_text_file_info(filename, file_size_mb)
                return self._get_text_load_message(filename)
                
            elif file_extension == 'docx':
                self.text_content = self.extract_text_from_docx(filename)
                self._set_text_file_info(filename, file_size_mb)
                return self._get_text_load_message(filename)
                
            elif file_extension == 'pdf':
                self.text_content = self.extract_text_from_pdf(filename)
                self._set_text_file_info(filename, file_size_mb)
                return self._get_text_load_message(filename)
                
            elif file_extension == 'pptx':
                self.text_content = self.extract_text_from_pptx(filename)
                self._set_text_file_info(filename, file_size_mb)
                return self._get_text_load_message(filename)
            

            elif file_extension in ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif']:
                image_analysis = self.analyze_image_content(filename)
                self.text_content = image_analysis.get('extracted_text', '')
                self._set_image_file_info(filename, file_size_mb, image_analysis)
                return self._get_image_load_message(filename, image_analysis)
            
            else:
                return f" Unsupported file format: {file_extension}"
                
        except Exception as e:
            return f" Error loading file: {str(e)}"
    
    def _set_structured_file_info(self, filename: str, file_size_mb: float):
        """Set file info for structured data"""
        self.file_info = {
            'filename': filename,
            'file_type': 'structured',
            'file_size_mb': file_size_mb,
            'shape': self.dataset.shape,
            'columns': list(self.dataset.columns),
            'dtypes': self.dataset.dtypes.to_dict(),
            'loaded_at': datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        }
    
    def _set_text_file_info(self, filename: str, file_size_mb: float):
        """Set file info for text content"""
        words = len(self.text_content.split())
        chars = len(self.text_content)
        lines = len(self.text_content.split('\n'))
        
        self.file_info = {
            'filename': filename,
            'file_type': 'text',
            'file_size_mb': file_size_mb,
            'word_count': words,
            'character_count': chars,
            'line_count': lines,
            'loaded_at': datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        }
    
    def _set_image_file_info(self, filename: str, file_size_mb: float, image_analysis: Dict):
        """Set file info for images"""
        self.file_info = {
            'filename': filename,
            'file_type': 'image',
            'file_size_mb': file_size_mb,
            'image_format': image_analysis.get('format'),
            'dimensions': image_analysis.get('size'),
            'has_text': image_analysis.get('has_text', False),
            'loaded_at': datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        }
        self.file_info.update(image_analysis)
    
    def _get_structured_load_message(self, filename: str) -> str:
        """Get load message for structured data"""
        return f" Successfully loaded structured data '{filename}'\n    Shape: {self.dataset.shape[0]} rows × {self.dataset.shape[1]} columns\n    Columns: {', '.join(self.dataset.columns[:5])}{'...' if len(self.dataset.columns) > 5 else ''}\n    Size: {self.file_info['file_size_mb']:.2f} MB"
    
    def _get_text_load_message(self, filename: str) -> str:
        """Get load message for text content"""
        return f" Successfully loaded text document '{filename}'\n    Words: {self.file_info['word_count']:,}\n    Lines: {self.file_info['line_count']:,}\n    Size: {self.file_info['file_size_mb']:.2f} MB"
    
    def _get_image_load_message(self, filename: str, image_analysis: Dict) -> str:
        """Get load message for images"""
        dimensions = image_analysis.get('size', (0, 0))
        has_text = image_analysis.get('has_text', False)
        return f" Successfully loaded image '{filename}'\n    Dimensions: {dimensions[0]}×{dimensions[1]} pixels\n   📝 Contains text: {'Yes' if has_text else 'No'}\n    Size: {self.file_info['file_size_mb']:.2f} MB"
    
    def get_file_summary(self) -> str:
        """Get comprehensive file summary"""
        if not self.file_info:
            return " No file loaded. Use 'load <filename>' to load a file first."
        
        try:
            summary = []
            file_type = self.file_info.get('file_type', 'unknown')
            

            summary.append(" **File Overview**")
            summary.append(f"   • File: {self.file_info.get('filename', 'Unknown')}")
            summary.append(f"   • Type: {file_type.title()}")
            summary.append(f"   • Size: {self.file_info.get('file_size_mb', 0):.2f} MB")
            summary.append(f"   • Loaded: {self.file_info.get('loaded_at', 'Unknown')}")
            
            if file_type == 'structured':

                summary.append(f"\n **Data Structure**")
                summary.append(f"   • Shape: {self.dataset.shape[0]} rows × {self.dataset.shape[1]} columns")
                summary.append(f"   • Memory Usage: {self.dataset.memory_usage(deep=True).sum() / 1024:.2f} KB")
                
                summary.append(f"\n **Column Information**")
                for col in self.dataset.columns[:10]: 
                    dtype = str(self.dataset[col].dtype)
                    null_count = self.dataset[col].isnull().sum()
                    null_pct = (null_count / len(self.dataset)) * 100
                    summary.append(f"   • {col}: {dtype} ({null_count} nulls, {null_pct:.1f}%)")
                
                if len(self.dataset.columns) > 10:
                    summary.append(f"   ... and {len(self.dataset.columns) - 10} more columns")
                

                missing_data = self.dataset.isnull().sum()
                if missing_data.sum() > 0:
                    summary.append(f"\n **Missing Values**")
                    for col, missing in missing_data[missing_data > 0].head(5).items():
                        pct = (missing / len(self.dataset)) * 100
                        summary.append(f"   • {col}: {missing} ({pct:.1f}%)")
            
            elif file_type == 'text':

                summary.append(f"\n **Content Analysis**")
                summary.append(f"   • Words: {self.file_info.get('word_count', 0):,}")
                summary.append(f"   • Characters: {self.file_info.get('character_count', 0):,}")
                summary.append(f"   • Lines: {self.file_info.get('line_count', 0):,}")
                

                if self.text_content:
                    words = self.text_content.split()
                    unique_words = len(set(word.lower() for word in words if word.isalpha()))
                    avg_word_length = np.mean([len(word) for word in words if word.isalpha()]) if words else 0
                    summary.append(f"   • Unique words: {unique_words:,}")
                    summary.append(f"   • Average word length: {avg_word_length:.1f} characters")
            
            elif file_type == 'image':

                summary.append(f"\n **Image Properties**")
                summary.append(f"   • Format: {self.file_info.get('image_format', 'Unknown')}")
                summary.append(f"   • Dimensions: {self.file_info.get('width', 0)}×{self.file_info.get('height', 0)} pixels")
                summary.append(f"   • Mode: {self.file_info.get('mode', 'Unknown')}")
                summary.append(f"   • Has transparency: {self.file_info.get('has_transparency', False)}")
                summary.append(f"   • Contains text: {self.file_info.get('has_text', False)}")
            
            return "\n".join(summary)
            
        except Exception as e:
            return f" Error generating summary: {str(e)}"
    
    def analyze_content(self) -> str:
        """Perform comprehensive content analysis"""
        if not self.file_info:
            return " No file loaded. Use 'load <filename>' to load a file first."
        
        try:
            basic_info = self.get_file_summary()
            file_type = self.file_info.get('file_type', 'unknown')
            

            if file_type == 'structured':
                context = f"""
                Dataset Analysis Context:
                - Shape: {self.dataset.shape}
                - Columns: {', '.join(self.dataset.columns)}
                - Data types: {self.dataset.dtypes.to_dict()}
                - Missing values: {self.dataset.isnull().sum().to_dict()}
                
                Sample data (first 5 rows):
                {self.dataset.head().to_string()}
                
                Statistical summary:
                {self.dataset.describe().to_string()}
                """
                
                prompt = """
                Analyze this structured dataset and provide:
                1. Key insights about the data structure and quality
                2. Patterns, trends, or correlations you notice
                3. Data quality issues and recommendations
                4. Suggested visualizations and analysis approaches
                5. Business insights or interesting observations
                6. Recommended next steps for deeper analysis
                
                Be specific and actionable in your recommendations.
                """
                
            elif file_type == 'text':

                text_analysis = self._perform_text_analysis()
                context = f"""
                Text Document Analysis Context:
                - Word count: {self.file_info.get('word_count', 0)}
                - Character count: {self.file_info.get('character_count', 0)}
                - Text analysis results: {text_analysis}
                
                Sample text (first 1000 characters):
                {self.text_content[:1000]}...
                """
                
                prompt = """
                Analyze this text document and provide:
                1. Content summary and main themes
                2. Sentiment analysis and tone assessment
                3. Key topics and concepts identified
                4. Document structure and organization
                5. Language characteristics and style
                6. Insights and recommendations for further analysis
                
                Focus on actionable insights and interesting patterns.
                """
                
            elif file_type == 'image':
                context = f"""
                Image Analysis Context:
                - Format: {self.file_info.get('image_format')}
                - Dimensions: {self.file_info.get('width')}x{self.file_info.get('height')}
                - Contains text: {self.file_info.get('has_text')}
                - Extracted text: {self.text_content[:500] if self.text_content else 'None'}
                """
                
                prompt = """
                Analyze this image file and provide:
                1. Technical image properties assessment
                2. Content analysis (if text was extracted)
                3. Potential use cases and applications
                4. Quality assessment and recommendations
                5. Suggestions for further processing or analysis
                
                Focus on practical insights and recommendations.
                """
            
            else:
                context = "Unknown file type"
                prompt = "Provide general file analysis based on available information."
            

            llm_analysis = self.query_llm(prompt, context)
            

            full_analysis = f"{basic_info}\n\n **AI Analysis Insights**\n{llm_analysis}"
            

            self.analysis_history.append({
                'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                'type': 'comprehensive_analysis',
                'file_type': file_type,
                'result': full_analysis
            })
            
            return full_analysis
            
        except Exception as e:
            return f" Error during analysis: {str(e)}"
        
        
    def _perform_text_analysis(self) -> Dict[str, Any]:
        """Perform detailed text analysis"""
        
        if not self.text_content:
            return {}

            analysis = {}


            words = self.text_content.split()
            analysis['word_count'] = len(words)
            analysis['unique_words'] = len(set(word.lower() for word in words if word.isalpha()))
            analysis['sentences'] = len([s for s in self.text_content.split('.') if s.strip()])


        try:
            from nltk.sentiment import SentimentIntensityAnalyzer
            import nltk
            nltk.download('vader_lexicon', quiet=True)
            sia = SentimentIntensityAnalyzer()
            sentiment_scores = sia.polarity_scores(self.text_content)
            analysis['sentiment'] = sentiment_scores
        except Exception as e:
            analysis['sentiment'] = f"Sentiment analysis unavailable: {str(e)}"


        try:
            from nltk.corpus import stopwords
            import nltk
            nltk.download('stopwords', quiet=True)
            stop_words = set(stopwords.words('english'))
            words_clean = [word.lower() for word in words if word.isalpha() and word.lower() not in stop_words]
            word_freq = pd.Series(words_clean).value_counts().head(10).to_dict()
            analysis['top_words'] = word_freq
        except Exception as e:

            words_clean = [word.lower() for word in words if word.isalpha() and len(word) > 3]
            word_freq = pd.Series(words_clean).value_counts().head(10).to_dict()
            analysis['top_words'] = word_freq
            analysis['stopwords_note'] = "Used basic filtering (NLTK stopwords unavailable)"

            return analysis

        except Exception as e:
            return {'error': str(e)}
    
    def answer_question(self, question: str) -> str:
        """Answer specific questions about the loaded content"""
        if not self.file_info:
            return " No file loaded. Use 'load <filename>' to load a file first."
        
        try:
            file_type = self.file_info.get('file_type', 'unknown')
            

            if file_type == 'structured' and self.dataset is not None:
                context = f"""
                Current Dataset Information:
                - Shape: {self.dataset.shape}
                - Columns: {', '.join(self.dataset.columns)}
                - Data types: {self.dataset.dtypes.to_dict()}
                
                Sample data (first 10 rows):
                {self.dataset.head(10).to_string()}
                
                Statistical summary:
                {self.dataset.describe().to_string()}
                """
                
            elif file_type in ['text', 'image'] and self.text_content:
                context = f"""
                Content Information:
                - File type: {file_type}
                - Content length: {len(self.text_content)} characters
                
                Content sample:
                {self.text_content[:2000]}
                """
                
            else:
                context = f"File information: {self.file_info}"
            
            prompt = f"""
            Based on the {file_type} content provided, please answer this question: {question}
            
            Please provide:
            1. A direct answer to the question
            2. Supporting evidence from the data/content
            3. Relevant statistics, calculations, or quotes
            4. Additional insights or related observations
            5. Suggestions for follow-up questions or analysis
            
            Be specific and use actual data/content where possible.
            """
            
            response = self.query_llm(prompt, context)
            

            self.analysis_history.append({
                'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                'type': 'question_answer',
                'question': question,
                'answer': response,
                'file_type': file_type
            })
            
            return f" **Question:** {question}\n\n **Answer:**\n{response}"
            
        except Exception as e:
            return f" Error answering question: {str(e)}"
    
    def create_visualization(self, viz_type: str = "auto", column: str = None, **kwargs) -> str:
        """Create visualizations for the loaded data"""
        if self.file_info.get('file_type') != 'structured' or self.dataset is None:
            return " Visualizations are only available for structured data (CSV, Excel, JSON arrays)."
        
        try:
            plt.style.use('default')
            fig, axes = plt.subplots(2, 2, figsize=(15, 12))
            fig.suptitle(f'Data Analysis Dashboard - {self.file_info.get("filename", "Dataset")}', fontsize=16)
            

            ax1 = axes[0, 0]
            numeric_cols = self.dataset.select_dtypes(include=[np.number]).columns
            if len(numeric_cols) > 0:
                self.dataset[numeric_cols[:5]].hist(ax=ax1, bins=20, alpha=0.7)
                ax1.set_title('Numeric Columns Distribution')
            else:
                ax1.text(0.5, 0.5, 'No numeric columns found', ha='center', va='center')
                ax1.set_title('Data Distribution')
            

            ax2 = axes[0, 1]
            missing_data = self.dataset.isnull()
            if missing_data.sum().sum() > 0:
                sns.heatmap(missing_data.iloc[:50], ax=ax2, cbar=True, cmap='viridis')
                ax2.set_title('Missing Values Pattern (First 50 rows)')
            else:
                ax2.text(0.5, 0.5, 'No missing values', ha='center', va='center')
                ax2.set_title('Missing Values Analysis')
            

            ax3 = axes[1, 0]
            if len(numeric_cols) > 1:
                corr_matrix = self.dataset[numeric_cols].corr()
                sns.heatmap(corr_matrix, annot=True, ax=ax3, cmap='coolwarm', center=0)
                ax3.set_title('Correlation Matrix')
            else:
                ax3.text(0.5, 0.5, 'Insufficient numeric data\nfor correlation analysis', ha='center', va='center')
                ax3.set_title('Correlation Analysis')
            

            ax4 = axes[1, 1]
            null_counts = self.dataset.isnull().sum()
            null_counts = null_counts[null_counts > 0]
            if len(null_counts) > 0:
                null_counts.plot(kind='bar', ax=ax4)
                ax4.set_title('Null Values by Column')
                ax4.tick_params(axis='x', rotation=45)
            else:
                data_types = self.dataset.dtypes.value_counts()
                data_types.plot(kind='pie', ax=ax4, autopct='%1.1f%%')
                ax4.set_title('Data Types Distribution')
            
            plt.tight_layout()
            

            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            viz_filename = f"visualization_{timestamp}.png"
            plt.savefig(viz_filename, dpi=300, bbox_inches='tight')
            plt.show()
            

            if viz_type == "interactive":
                self._create_interactive_plots()
            
            viz_info = {
                'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                'type': 'data_dashboard',
                'filename': viz_filename,
                'description': 'Comprehensive data analysis dashboard'
            }
            
            self.visualizations.append(viz_info)
            
            return f" Visualization created successfully!\n    Saved as: {viz_filename}\n    Type: Data Analysis Dashboard\n    Focus: Overview, distributions, correlations, missing values"
            
        except Exception as e:
            return f" Error creating visualization: {str(e)}"
    
    def _create_interactive_plots(self):
        """Create interactive Plotly visualizations"""
        try:
            numeric_cols = self.dataset.select_dtypes(include=[np.number]).columns
            
            if len(numeric_cols) >= 2:

                fig = px.scatter(self.dataset, x=numeric_cols[0], y=numeric_cols[1], 
                               title=f'Interactive Scatter Plot: {numeric_cols[0]} vs {numeric_cols[1]}')
                fig.show()
                

                fig2 = px.box(self.dataset, y=numeric_cols[0], 
                             title=f'Box Plot: {numeric_cols[0]}')
                fig2.show()
                
        except Exception as e:
            print(f"Error creating interactive plots: {str(e)}")
    
    def export_analysis(self, format_type: str = "json") -> str:
        """Export analysis results and history"""
        try:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            
            export_data = {
                'file_info': self.file_info,
                'analysis_history': self.analysis_history,
                'visualizations': self.visualizations,
                'export_timestamp': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                'total_analyses': len(self.analysis_history)
            }
            
            if format_type.lower() == "json":
                filename = f"analysis_export_{timestamp}.json"
                with open(filename, 'w', encoding='utf-8') as f:
                    json.dump(export_data, f, indent=2, default=str)
                    
            elif format_type.lower() == "txt":
                filename = f"analysis_export_{timestamp}.txt"
                with open(filename, 'w', encoding='utf-8') as f:
                    f.write("UNIVERSAL DATA ANALYST - ANALYSIS EXPORT\n")
                    f.write("=" * 50 + "\n\n")
                    
                    f.write("FILE INFORMATION:\n")
                    f.write("-" * 20 + "\n")
                    for key, value in self.file_info.items():
                        f.write(f"{key}: {value}\n")
                    
                    f.write(f"\nANALYSIS HISTORY ({len(self.analysis_history)} entries):\n")
                    f.write("-" * 30 + "\n")
                    for i, analysis in enumerate(self.analysis_history, 1):
                        f.write(f"\n{i}. {analysis.get('type', 'Unknown').upper()}\n")
                        f.write(f"   Timestamp: {analysis.get('timestamp', 'Unknown')}\n")
                        if 'question' in analysis:
                            f.write(f"   Question: {analysis['question']}\n")
                        f.write(f"   Result: {str(analysis.get('result', analysis.get('answer', 'N/A')))[:500]}...\n")
                    
                    f.write(f"\nVISUALIZATIONS ({len(self.visualizations)} created):\n")
                    f.write("-" * 25 + "\n")
                    for viz in self.visualizations:
                        f.write(f"- {viz.get('description', 'Unknown')} ({viz.get('timestamp', 'Unknown')})\n")
            
            return f" Analysis exported successfully!\n    File: {filename}\n    Format: {format_type.upper()}\n    Analyses: {len(self.analysis_history)}\n    Visualizations: {len(self.visualizations)}"
            
        except Exception as e:
            return f" Error exporting analysis: {str(e)}"
    
    def generate_report(self) -> str:
        """Generate a comprehensive analysis report"""
        if not self.file_info:
            return " No file loaded. Use 'load <filename>' to load a file first."
        
        try:

            file_summary = self.get_file_summary()
            

            file_type = self.file_info.get('file_type', 'unknown')
            
            if file_type == 'structured':
                context = f"""
                Dataset Report Context:
                {file_summary}
                
                Dataset shape: {self.dataset.shape}
                Columns: {list(self.dataset.columns)}
                Data sample:
                {self.dataset.head().to_string()}
                
                Statistical summary:
                {self.dataset.describe().to_string()}
                
                Analysis history: {len(self.analysis_history)} previous analyses
                """
                
                prompt = """
                Generate a comprehensive data analysis report with the following sections:
                
                ## Executive Summary
                - Key findings and insights
                - Data quality assessment
                - Main recommendations
                
                ## Dataset Overview
                - Structure and composition
                - Data types and formats
                - Size and scope analysis
                
                ## Data Quality Analysis
                - Missing values assessment
                - Outliers and anomalies
                - Data consistency issues
                
                ## Statistical Insights
                - Descriptive statistics summary
                - Key patterns and trends
                - Correlations and relationships
                
                ## Business Insights
                - Actionable insights
                - Potential use cases
                - Value propositions
                
                ## Recommendations
                - Data cleaning suggestions
                - Further analysis opportunities
                - Visualization recommendations
                
                ## Technical Notes
                - Processing considerations
                - Limitations and assumptions
                - Next steps
                
                Make it professional, detailed, and actionable.
                """
                
            elif file_type == 'text':
                text_analysis = self._perform_text_analysis()
                context = f"""
                Text Document Report Context:
                {file_summary}
                
                Text analysis results: {text_analysis}
                Content sample: {self.text_content[:1000]}
                Analysis history: {len(self.analysis_history)} previous analyses
                """
                
                prompt = """
                Generate a comprehensive text analysis report with the following sections:
                
                ## Executive Summary
                - Document overview and purpose
                - Key themes and topics
                - Main insights and findings
                
                ## Content Analysis
                - Document structure and organization
                - Writing style and tone
                - Language characteristics
                
                ## Thematic Analysis
                - Major topics and concepts
                - Sentiment and emotional tone
                - Key messages and arguments
                
                ## Statistical Overview
                - Word count and vocabulary analysis
                - Readability assessment
                - Content distribution
                
                ## Insights and Patterns
                - Notable patterns or trends
                - Interesting observations
                - Content quality assessment
                
                ## Recommendations
                - Further analysis suggestions
                - Potential applications
                - Content improvement opportunities
                
                ## Technical Notes
                - Processing methodology
                - Limitations and considerations
                - Suggested tools for deeper analysis
                
                Make it comprehensive and insightful.
                """
                
            elif file_type == 'image':
                context = f"""
                Image Analysis Report Context:
                {file_summary}
                
                Image properties: {self.file_info}
                Extracted text: {self.text_content[:500] if self.text_content else 'None'}
                Analysis history: {len(self.analysis_history)} previous analyses
                """
                
                prompt = """
                Generate a comprehensive image analysis report with the following sections:
                
                ## Executive Summary
                - Image overview and characteristics
                - Technical assessment
                - Key findings
                
                ## Technical Analysis
                - Format and specifications
                - Resolution and quality assessment
                - File size and optimization
                
                ## Content Analysis
                - Visual content description
                - Text extraction results
                - Potential use cases
                
                ## Quality Assessment
                - Image quality evaluation
                - Technical recommendations
                - Processing suggestions
                
                ## Applications and Use Cases
                - Potential applications
                - Recommended processing workflows
                - Integration possibilities
                
                ## Recommendations
                - Optimization suggestions
                - Further analysis opportunities
                - Best practices for usage
                
                ## Technical Notes
                - Processing considerations
                - Tool recommendations
                - Limitations and assumptions
                
                Make it thorough and practical.
                """
            
            else:
                context = f"File information: {self.file_info}"
                prompt = "Generate a general file analysis report based on available information."
            

            report_content = self.query_llm(prompt, context)
            

            timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            report_header = f"""
# COMPREHENSIVE ANALYSIS REPORT
**Generated by Universal Data Analyst Agent**
**Date:** {timestamp}
**File:** {self.file_info.get('filename', 'Unknown')}
**Type:** {file_type.title()}

---

"""
            
            full_report = report_header + report_content
            

            report_filename = f"analysis_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md"
            with open(report_filename, 'w', encoding='utf-8') as f:
                f.write(full_report)
            

            self.analysis_history.append({
                'timestamp': timestamp,
                'type': 'comprehensive_report',
                'file_type': file_type,
                'report_file': report_filename,
                'result': 'Comprehensive report generated'
            })
            
            return f" Comprehensive report generated!\n    File: {report_filename}\n    Type: {file_type.title()} Analysis Report\n    Sections: Executive Summary, Analysis, Insights, Recommendations\n\n{full_report[:1000]}..."
            
        except Exception as e:
            return f" Error generating report: {str(e)}"
    
    def get_analysis_history(self) -> str:
        """Get formatted analysis history"""
        if not self.analysis_history:
            return " No analysis history available."
        
        try:
            history = [" **Analysis History**\n" + "=" * 30]
            
            for i, analysis in enumerate(self.analysis_history, 1):
                timestamp = analysis.get('timestamp', 'Unknown')
                analysis_type = analysis.get('type', 'Unknown').replace('_', ' ').title()
                file_type = analysis.get('file_type', 'unknown')
                
                history.append(f"\n{i}. **{analysis_type}** ({file_type})")
                history.append(f"    {timestamp}")
                
                if 'question' in analysis:
                    history.append(f"    Question: {analysis['question'][:100]}...")
                
                if 'result' in analysis:
                    result_preview = str(analysis['result'])[:200].replace('\n', ' ')
                    history.append(f"    Result: {result_preview}...")
                elif 'answer' in analysis:
                    answer_preview = str(analysis['answer'])[:200].replace('\n', ' ')
                    history.append(f"    Answer: {answer_preview}...")
            
            history.append(f"\n **Summary:** {len(self.analysis_history)} total analyses performed")
            
            return "\n".join(history)
            
        except Exception as e:
            return f" Error retrieving history: {str(e)}"
    
    def help_menu(self) -> str:
        """Display comprehensive help menu"""
        return """
 **UNIVERSAL DATA ANALYST AGENT - HELP MENU**
=====================================================

 **FILE OPERATIONS**
• load <filename>          - Load any file type (CSV, Excel, PDF, Word, PowerPoint, Images, Text)
• summary                  - Get comprehensive file summary and overview
• info                     - Display current file information

 **ANALYSIS COMMANDS**
• analyze                  - Perform comprehensive AI-powered content analysis
• ask <question>           - Ask specific questions about your data/content
• report                   - Generate detailed analysis report

 **VISUALIZATIONS** (Structured Data Only)
• visualize                - Create data analysis dashboard
• visualize interactive    - Create interactive Plotly visualizations
• chart <type> <column>    - Create specific chart types

 **EXPORT & HISTORY**
• export json/txt          - Export analysis results and history
• history                  - View all previous analyses
• visualizations           - List created visualizations

 **SUPPORTED FILE TYPES**
• **Structured Data:** CSV, Excel (.xlsx, .xls), JSON arrays
• **Documents:** PDF, Word (.docx), PowerPoint (.pptx), Text (.txt)
• **Images:** JPG, PNG, BMP, TIFF, GIF (with OCR text extraction)

 **AI CAPABILITIES**
• Advanced pattern recognition and insights
• Natural language querying of your data
• Automated data quality assessment
• Intelligent visualization recommendations
• Multi-format content understanding
• Sentiment analysis for text content
• Statistical analysis and correlations

 **EXAMPLE COMMANDS**
• load sales_data.csv
• analyze
• ask "What are the top selling products?"
• visualize interactive
• report
• export json

 **GETTING STARTED**
1. Load your file: load <filename>
2. Get overview: summary
3. Analyze content: analyze
4. Ask questions: ask <your question>
5. Create visuals: visualize
6. Generate report: report

Type any command to get started! 
"""

def main():
    """Main interactive loop for the Universal Data Analyst Agent"""
    print(" UNIVERSAL DATA ANALYST AGENT")
    print("=" * 50)
    print("AI-Powered Analysis for ALL File Types!")
    print("Supports: CSV, Excel, PDF, Word, PowerPoint, Images, Text files")
    print("Type 'help' for commands or 'quit' to exit\n")
    
# API KEY
    api_key = input("Enter your Together.ai API key: ").strip()
    if not api_key:
        print(" API key required to continue.")
        return
    
    agent = UniversalDataAnalystAgent(api_key)
    
    print(f"\n Agent initialized successfully!")
    print("Ready to analyze your data! 🚀\n")
    
    print(agent.help_menu())
    
    while True:
        try:
            command = input("🔍 Enter command: ").strip().lower()
            
            if command in ['quit', 'exit', 'q']:
                print("👋 Thank you for using Universal Data Analyst Agent!")
                break
            
            elif command == 'help':
                print(agent.help_menu())
            
            elif command.startswith('load '):
                filename = command[5:].strip()
                result = agent.load_file(filename)
                print(result)
            
            elif command in ['summary', 'info']:
                result = agent.get_file_summary()
                print(result)
            
            elif command == 'analyze':
                print(" Performing comprehensive analysis...")
                result = agent.analyze_content()
                print(result)
            
            elif command.startswith('ask '):
                question = command[4:].strip()
                print(" Analyzing your question...")
                result = agent.answer_question(question)
                print(result)
            
            elif command in ['visualize', 'viz']:
                print(" Creating visualizations...")
                result = agent.create_visualization()
                print(result)
            
            elif command == 'visualize interactive':
                print(" Creating interactive visualizations...")
                result = agent.create_visualization(viz_type="interactive")
                print(result)
            
            elif command == 'report':
                print(" Generating comprehensive report...")
                result = agent.generate_report()
                print(result)
            
            elif command.startswith('export '):
                format_type = command[7:].strip()
                result = agent.export_analysis(format_type)
                print(result)
            
            elif command == 'history':
                result = agent.get_analysis_history()
                print(result)
            
            elif command == '':
                continue
            
            else:
                print(" Unknown command. Type 'help' for available commands.")
        
        except KeyboardInterrupt:
            print("\n👋 Goodbye!")
            break
        except Exception as e:
            print(f" An error occurred: {str(e)}")
            print("Please try again or type 'help' for assistance.")

if __name__ == "__main__":
    main()